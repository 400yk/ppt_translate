import os
import tempfile
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from services.llm_service import gemini_batch_translate_with_size
from utils.pptx_utils import measure_text_bbox, fit_font_size_to_bbox, fit_font_size_for_title
from config import (
    DEFAULT_FONT_NAME, DEFAULT_FONT_SIZE, DEFAULT_TITLE_FONT_SIZE
)

def apply_font_color(run, font_color=None, font_color_type=None, theme_color=None):
    """
    Apply font color to a run, handling both RGB and theme colors.
    
    Args:
        run: The text run to apply color to
        font_color: RGB color object (if available)
        font_color_type: Type of color ('rgb' or 'theme')
        theme_color: Theme color enum (if available)
    """
    try:
        if font_color_type == 'rgb' and font_color is not None:
            # Apply RGB color
            run.font.color.rgb = font_color
            print(f"Applied RGB color: {font_color}")
        elif font_color_type == 'theme' and theme_color is not None:
            # Apply theme color
            run.font.color.theme_color = theme_color
            print(f"Applied theme color: {theme_color}")
        else:
            # No color information available - leave as default
            # This preserves any inherited theme colors
            print("No explicit color found - using inherited/default color")
    except Exception as e:
        print(f"Warning: Could not apply font color: {e}")
        # If color application fails, continue without color to avoid breaking the translation

def register_routes(app):
    """
    Register translation routes with the Flask app.
    This function is now just a wrapper for the translate_api Blueprint.
    Routes are now defined in api/translate_api.py.
    We keep this function for backward compatibility.
    """
    pass

def translate_pptx(input_stream, src_lang, dest_lang):
    """
    Translate a PowerPoint file from source language to destination language.
    
    Args:
        input_stream: File-like object containing the PowerPoint file
        src_lang: Source language code (e.g. 'zh' for Chinese)
        dest_lang: Destination language code (e.g. 'en' for English)
        
    Returns:
        Tuple containing:
        - path to the translated file
        - total character count that was translated
    """
    prs = Presentation(input_stream)
    text_shapes = []
    texts = []
    table_cells = []
    table_texts = []
    
    # Collect all text shapes and their texts
    for slide in prs.slides:
        for shape in slide.shapes:
            # Handle regular text shapes
            if hasattr(shape, "text") and shape.text.strip():
                text_shapes.append(shape)
                texts.append(shape.text)
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            table_cells.append(cell)
                            table_texts.append(cell.text)

    # Batch translate all text content together using the new batched approach
    all_texts = texts + table_texts
    print(f"Total texts to translate: {len(all_texts)}")
    print(f"Translating from {src_lang} to {dest_lang}")
    all_translated_texts, total_characters = gemini_batch_translate_with_size(all_texts, src_lang, dest_lang, batch_size=200)
    
    # Split the translated texts back into shape texts and table texts
    translated_texts = all_translated_texts[:len(texts)]
    translated_table_texts = all_translated_texts[len(texts):]
    
    print(f"Collected {len(texts)} texts from regular shapes.")
    print(f"Collected {len(table_texts)} texts from table cells.")
    print(f"Received {len(all_translated_texts)} translated texts total.")
    print(f"Sample original: {all_texts[:3]}")
    print(f"Sample translated: {all_translated_texts[:3]}")
    print(f"Total characters: {total_characters}")
    print(f"Starting text formatting preservation process...")
    
    # Update regular text shapes
    for shape, translated in zip(text_shapes, translated_texts):
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            text_frame = shape.text_frame
            # Get original font properties from the first run (if available)
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                original_run = text_frame.paragraphs[0].runs[0]
                font_name = original_run.font.name or DEFAULT_FONT_NAME
                
                # Style scale factor for text formatting (bold/italic/underline)
                style_scale_factor = 1.0
                
                # Check if text has formatting that might need more space
                if original_run.font.bold:
                    style_scale_factor *= 0.9  # Bold text needs ~10% more space
                if original_run.font.italic:
                    style_scale_factor *= 0.95  # Italic text needs ~5% more space
                if original_run.font.underline:
                    style_scale_factor *= 0.95  # Underlined text needs ~5% more space
                
                # Check if this is a title placeholder
                is_title = False
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Title placeholders have types: TITLE (1) or CENTER_TITLE (3)
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        is_title = True
                
                # Use appropriate default font size based on whether it's a title
                if original_run.font.size:
                    original_font_size = int(original_run.font.size.pt)
                else:
                    original_font_size = DEFAULT_TITLE_FONT_SIZE if is_title else DEFAULT_FONT_SIZE
                
                # Enhanced font color extraction: run -> paragraph -> text_frame
                font_color = None
                font_color_type = None
                theme_color = None
                
                # 1. Try run-level color
                if original_run.font.color is not None:
                    color_obj = original_run.font.color
                    try:
                        # First try RGB color
                        font_color = color_obj.rgb
                        font_color_type = 'rgb'
                    except AttributeError:
                        try:
                            # Try theme color
                            if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                                theme_color = color_obj.theme_color
                                font_color_type = 'theme'
                        except AttributeError:
                            pass
                
                # 2. Try paragraph-level color if we haven't found one yet
                if not font_color and not theme_color:
                    para = text_frame.paragraphs[0]
                    if para.font and para.font.color is not None:
                        color_obj = para.font.color
                        try:
                            font_color = color_obj.rgb
                            font_color_type = 'rgb'
                        except AttributeError:
                            try:
                                if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                                    theme_color = color_obj.theme_color
                                    font_color_type = 'theme'
                            except AttributeError:
                                pass
            else:
                font_name = DEFAULT_FONT_NAME
                style_scale_factor = 1.0  # Default for no formatting
                
                # Check if this is a title placeholder
                is_title = False
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        is_title = True
                
                original_font_size = DEFAULT_TITLE_FONT_SIZE if is_title else DEFAULT_FONT_SIZE
                font_color = None
                
            # Measure the bounding box of the original text
            original_text = shape.text
            orig_w, orig_h = measure_text_bbox(original_text, font_name, original_font_size)
            
            # For titles, only constrain height, not width
            if is_title:
                best_font_size = fit_font_size_for_title(orig_h, translated, font_name, original_font_size)
            else:
                # For regular content, constrain both dimensions
                best_font_size = fit_font_size_to_bbox(orig_w, orig_h, translated, font_name, original_font_size)
            
            # Apply the style scaling factor to the font size
            best_font_size = int(best_font_size * style_scale_factor)
                
            # Store original formatting information before clearing
            original_paragraphs = []
            for para in text_frame.paragraphs:
                para_info = {
                    'level': para.level,
                    'alignment': para.alignment,
                    'runs': []
                }
                
                # Store runs and their formatting
                for run in para.runs:
                    run_info = {
                        'text_length': len(run.text),
                        'bold': run.font.bold,
                        'italic': run.font.italic, 
                        'underline': run.font.underline,
                        'font_name': run.font.name or font_name,
                        'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') else None
                    }
                    para_info['runs'].append(run_info)
                
                original_paragraphs.append(para_info)
                
            text_frame.clear()
            
            # Simple approach for short texts - use a single formatted run
            if len(original_paragraphs) == 1 and len(original_paragraphs[0]['runs']) == 1:
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = translated
                run.font.name = font_name
                run.font.size = Pt(best_font_size)
                
                # Apply original formatting
                original_run_info = original_paragraphs[0]['runs'][0]
                if original_run_info['bold'] is not None:
                    run.font.bold = original_run_info['bold']
                if original_run_info['italic'] is not None:
                    run.font.italic = original_run_info['italic']
                if original_run_info['underline'] is not None:
                    run.font.underline = original_run_info['underline']
                
                # Apply color using the improved color handling
                apply_font_color(run, font_color, font_color_type, theme_color)
            else:
                # For more complex text, preserve the structure using paragraph count
                # This is a simplification as we can't perfectly map formatting from source to translated text
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = translated
                run.font.name = font_name
                run.font.size = Pt(best_font_size)
                
                # Apply formatting from dominant formatting in original text
                has_bold = any(run_info['bold'] for para in original_paragraphs for run_info in para['runs'] if run_info['bold'] is not None)
                has_italic = any(run_info['italic'] for para in original_paragraphs for run_info in para['runs'] if run_info['italic'] is not None)
                has_underline = any(run_info['underline'] for para in original_paragraphs for run_info in para['runs'] if run_info['underline'] is not None)
                
                run.font.bold = has_bold
                run.font.italic = has_italic
                run.font.underline = has_underline
                
                # Apply color using the improved color handling
                apply_font_color(run, font_color, font_color_type, theme_color)
        else:
            shape.text = translated
    
    # Update table cells with translated text
    for cell, translated in zip(table_cells, translated_table_texts):
        text_frame = cell.text_frame
        
        # Store original formatting if available
        original_paragraphs = []
        font_name = DEFAULT_FONT_NAME
        font_color = None
        original_font_size = DEFAULT_FONT_SIZE
        style_scale_factor = 1.0
        
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            original_run = text_frame.paragraphs[0].runs[0]
            font_name = original_run.font.name or DEFAULT_FONT_NAME
            
            # Style scale factor for text formatting (bold/italic/underline)
            style_scale_factor = 1.0
            
            # Check if text has formatting that might need more space
            if original_run.font.bold:
                style_scale_factor *= 0.9  # Bold text needs ~10% more space
            if original_run.font.italic:
                style_scale_factor *= 0.95  # Italic text needs ~5% more space
            if original_run.font.underline:
                style_scale_factor *= 0.95  # Underlined text needs ~5% more space
            
            # Get original font size
            if original_run.font.size:
                original_font_size = int(original_run.font.size.pt)
            
            # Get original font color with enhanced detection
            font_color_type = None
            theme_color = None
            if original_run.font.color is not None:
                color_obj = original_run.font.color
                try:
                    font_color = color_obj.rgb
                    font_color_type = 'rgb'
                except AttributeError:
                    try:
                        if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                            theme_color = color_obj.theme_color
                            font_color_type = 'theme'
                    except AttributeError:
                        pass
            
            # Store original formatting
            for para in text_frame.paragraphs:
                para_info = {
                    'level': para.level,
                    'alignment': para.alignment,
                    'runs': []
                }
                
                for run in para.runs:
                    run_info = {
                        'text_length': len(run.text),
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'underline': run.font.underline,
                        'font_name': run.font.name or font_name,
                        'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') else None
                    }
                    para_info['runs'].append(run_info)
                
                original_paragraphs.append(para_info)
        
        # Measure the original text's bounding box
        original_text = cell.text
        orig_w, orig_h = measure_text_bbox(original_text, font_name, original_font_size)
        
        # Get the best font size that fits
        best_font_size = fit_font_size_to_bbox(orig_w, orig_h, translated, font_name, original_font_size)
        best_font_size = int(best_font_size * style_scale_factor)
        
        # Clear and update the text frame
        text_frame.clear()
        
        # Simple case: just a single paragraph with a single run
        if len(original_paragraphs) == 1 and len(original_paragraphs[0]['runs']) == 1:
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = translated
            run.font.name = font_name
            run.font.size = Pt(best_font_size)
            
            # Apply original formatting
            original_run_info = original_paragraphs[0]['runs'][0]
            if original_run_info['bold'] is not None:
                run.font.bold = original_run_info['bold']
            if original_run_info['italic'] is not None:
                run.font.italic = original_run_info['italic']
            if original_run_info['underline'] is not None:
                run.font.underline = original_run_info['underline']
            
            # Apply color using the improved color handling
            apply_font_color(run, font_color, font_color_type, theme_color)
        else:
            # Complex case: use dominant formatting
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = translated
            run.font.name = font_name
            run.font.size = Pt(best_font_size)
            
            # Apply formatting from dominant formatting in original text
            has_bold = any(run_info['bold'] for para in original_paragraphs for run_info in para['runs'] if run_info['bold'] is not None)
            has_italic = any(run_info['italic'] for para in original_paragraphs for run_info in para['runs'] if run_info['italic'] is not None)
            has_underline = any(run_info['underline'] for para in original_paragraphs for run_info in para['runs'] if run_info['underline'] is not None)
            
            run.font.bold = has_bold
            run.font.italic = has_italic
            run.font.underline = has_underline
            
            # Apply color using the improved color handling
            apply_font_color(run, font_color, font_color_type, theme_color)
    
    # Save the translated presentation to a temporary file
    output_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(output_file)
    output_file.close()
    
    return output_file.name, total_characters 