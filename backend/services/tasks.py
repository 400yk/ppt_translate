import io
import os
from celery_init import celery_app # Import from celery_init
from services.translate_service import translate_pptx as original_translate_pptx
from services.s3_service import s3_service
from db.models import User, GuestTranslation, db # Assuming User and db are accessible
from celery.exceptions import Retry
import time

def calculate_translation_rate(original_texts, translated_texts):
    """
    Calculate the percentage of texts that were actually translated (not returned as original).
    
    Args:
        original_texts: List of original text strings
        translated_texts: List of translated text strings
    
    Returns:
        Float between 0.0 and 1.0 representing the translation success rate
    """
    if not original_texts or not translated_texts:
        return 0.0
    
    if len(original_texts) != len(translated_texts):
        return 0.0
    
    translated_count = sum(1 for orig, trans in zip(original_texts, translated_texts) if orig != trans)
    return translated_count / len(original_texts)

@celery_app.task(bind=True, autoretry_for=(Exception,), retry_kwargs={'max_retries': 2, 'countdown': 30}, retry_backoff=True)
def process_translation_task(self, user_id: int, original_file_content_bytes: bytes, original_filename: str, src_lang: str, dest_lang: str):
    """
    Celery task to translate a PPTX file.
    Now includes automatic retry logic for translation failures.
    
    Args:
        original_file_content_bytes: The content of the uploaded file as bytes.
        user_id: ID of the user requesting translation
        original_filename: Name of the original file
        src_lang: Source language code
        dest_lang: Destination language code
    """
    print(f"Celery task {self.request.id}: Starting translation for user {user_id}, file {original_filename} ({src_lang} → {dest_lang}) (attempt {self.request.retries + 1})")
    
    try:
        # Reconstruct a file-like stream from bytes
        file_stream = io.BytesIO(original_file_content_bytes)
        
        # Store original texts for comparison
        from pptx import Presentation
        temp_prs = Presentation(file_stream)
        original_texts = []
        
        # Collect all text content for rate calculation
        for slide in temp_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    original_texts.append(shape.text)
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                original_texts.append(cell.text)
        
        # Reset file stream
        file_stream.seek(0)
        
        # Call the original translation service function
        # This function saves the translated file to a temporary path on the worker
        translated_file_path, character_count = original_translate_pptx(file_stream, src_lang, dest_lang)
        
        # Check translation success by re-reading the translated file
        translated_prs = Presentation(translated_file_path)
        translated_texts = []
        
        # Collect translated text content
        for slide in translated_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    translated_texts.append(shape.text)
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                translated_texts.append(cell.text)
        
        # Calculate translation success rate
        translation_rate = calculate_translation_rate(original_texts, translated_texts)
        
        print(f"Celery task {self.request.id}: Translation rate: {translation_rate:.1%} ({int(translation_rate * len(original_texts))}/{len(original_texts)} texts translated) ({src_lang} → {dest_lang})")
        
        # If translation rate is very low (less than 10%), consider it a failure and retry
        if translation_rate < 0.1 and len(original_texts) > 0:
            error_msg = f"Translation failed ({src_lang} → {dest_lang}) - only {translation_rate:.1%} of texts were translated (likely due to API rate limiting or temporary issues)"
            print(f"Celery task {self.request.id}: {error_msg}")
            
            # Clean up the failed translation file
            try:
                os.unlink(translated_file_path)
            except Exception as e:
                print(f"Warning: Could not clean up failed translation file: {e}")
            
            # Calculate retry delay with shorter backoff
            retry_count = self.request.retries
            base_delay = 30  # Base delay of 30 seconds
            max_delay = 120  # Maximum delay of 2 minutes
            delay = min(base_delay * (2 ** retry_count), max_delay)
            
            print(f"Celery task {self.request.id}: Retrying ({src_lang} → {dest_lang}) in {delay} seconds (attempt {retry_count + 1}/2)")
            
            # Retry the task with exponential backoff
            raise self.retry(countdown=delay, exc=Exception(error_msg))
        
        # Fetch the user (ensure app context if not using ContextTask from make_celery)
        user = User.query.get(user_id)
        if not user:
            error_msg = f"User with ID {user_id} not found. Cannot record translation."
            print(f"Celery task {self.request.id}: {error_msg}")
            # Decide how to handle this - fail the task or just log?
            # For now, we'll let it proceed but the translation won't be recorded for the user.
            # raise FictionalTaskError(error_msg) # if you want to mark task as failed
        else:
            try:
                user.record_translation(original_filename, src_lang, dest_lang, character_count)
                db.session.commit() # Commit after recording
                print(f"Celery task {self.request.id}: Translation recorded for user {user_id}")
            except Exception as e:
                db.session.rollback()
                print(f"Celery task {self.request.id}: Error recording translation for user {user_id}: {e}")
                # Optionally re-raise to mark task as failed due to recording error
                # raise

        print(f"Celery task {self.request.id}: Translation successful ({src_lang} → {dest_lang}). File at: {translated_file_path}, Chars: {character_count}")
        
        # Try to upload to S3 with OSS fallback
        upload_result = {'success': False, 'service': None, 'key': None, 'download_url': None}
        
        if s3_service.is_available():
            try:
                upload_result = s3_service.upload_file_with_fallback(
                    translated_file_path, 
                    self.request.id, 
                    f"translated_{original_filename}"
                )
                
                if upload_result['success']:
                    service = upload_result['service']
                    key = upload_result['key']
                    download_url = upload_result['download_url']
                    
                    print(f"Celery task {self.request.id}: File uploaded to {service.upper()}: {key}")
                    
                    # Clean up local file after successful upload
                    try:
                        os.unlink(translated_file_path)
                        print(f"Celery task {self.request.id}: Local file cleaned up")
                    except Exception as e:
                        print(f"Celery task {self.request.id}: Warning - could not clean up local file: {e}")
                else:
                    print(f"Celery task {self.request.id}: Upload failed: {upload_result['error']}")
                    
            except Exception as e:
                print(f"Celery task {self.request.id}: Error during upload with fallback: {e}")
        else:
            print(f"Celery task {self.request.id}: S3 not available, using local storage")
        
        # The task needs to return information that the client can use.
        result = {
            'status': 'SUCCESS',
            'message': 'File translated successfully.',
            'original_filename': original_filename,
            'character_count': character_count,
            'translation_rate': translation_rate,
            'texts_translated': int(translation_rate * len(original_texts)),
            'total_texts': len(original_texts)
        }
        
        # Include appropriate file reference based on storage method
        if upload_result['success']:
            result['storage_key'] = upload_result['key']
            result['download_url'] = upload_result['download_url']
            result['storage_type'] = upload_result['service']
            # Keep backward compatibility
            if upload_result['service'] == 's3':
                result['s3_key'] = upload_result['key']
        else:
            result['translated_file_path'] = translated_file_path  # Path on the WORKER dyno
            result['storage_type'] = 'local'
        
        return result
    except Retry:
        # Re-raise Retry exceptions to let Celery handle them
        raise
    except Exception as e:
        print(f"Celery task {self.request.id}: Error during translation: {e}")
        import traceback
        traceback.print_exc()
        # Optionally, re-raise to mark the task as FAILED in Celery
        # For now, it returns an error structure
        # self.update_state(state='FAILURE', meta={'exc_type': type(e).__name__, 'exc_message': str(e)})
        raise # Re-raising will mark task as FAILED 

@celery_app.task(bind=True, autoretry_for=(Exception,), retry_kwargs={'max_retries': 2, 'countdown': 30}, retry_backoff=True)
def process_guest_translation_task(self, client_ip: str, original_file_content_bytes: bytes, original_filename: str, src_lang: str, dest_lang: str, estimated_character_count: int):
    """
    Celery task to translate a PPTX file for guest users.
    Now includes automatic retry logic for translation failures.
    """
    print(f"Celery guest task {self.request.id}: Starting translation for IP {client_ip}, file {original_filename} ({src_lang} → {dest_lang}) (attempt {self.request.retries + 1})")
    
    try:
        # Reconstruct a file-like stream from bytes
        file_stream = io.BytesIO(original_file_content_bytes)
        
        # Store original texts for comparison
        from pptx import Presentation
        temp_prs = Presentation(file_stream)
        original_texts = []
        
        # Collect all text content for rate calculation
        for slide in temp_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    original_texts.append(shape.text)
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                original_texts.append(cell.text)
        
        # Reset file stream
        file_stream.seek(0)
        
        # Call the original translation service function
        # This function saves the translated file to a temporary path on the worker
        translated_file_path, character_count = original_translate_pptx(file_stream, src_lang, dest_lang)
        
        # Check translation success by re-reading the translated file
        translated_prs = Presentation(translated_file_path)
        translated_texts = []
        
        # Collect translated text content
        for slide in translated_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    translated_texts.append(shape.text)
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                translated_texts.append(cell.text)
        
        # Calculate translation success rate
        translation_rate = calculate_translation_rate(original_texts, translated_texts)
        
        print(f"Celery guest task {self.request.id}: Translation rate: {translation_rate:.1%} ({int(translation_rate * len(original_texts))}/{len(original_texts)} texts translated) ({src_lang} → {dest_lang})")
        
        # If translation rate is very low (less than 10%), consider it a failure and retry
        if translation_rate < 0.1 and len(original_texts) > 0:
            error_msg = f"Translation failed ({src_lang} → {dest_lang}) - only {translation_rate:.1%} of texts were translated (likely due to API rate limiting or temporary issues)"
            print(f"Celery guest task {self.request.id}: {error_msg}")
            
            # Clean up the failed translation file
            try:
                os.unlink(translated_file_path)
            except Exception as e:
                print(f"Warning: Could not clean up failed translation file: {e}")
            
            # Calculate retry delay with shorter backoff
            retry_count = self.request.retries
            base_delay = 30  # Base delay of 30 seconds
            max_delay = 120  # Maximum delay of 2 minutes
            delay = min(base_delay * (2 ** retry_count), max_delay)
            
            print(f"Celery guest task {self.request.id}: Retrying ({src_lang} → {dest_lang}) in {delay} seconds (attempt {retry_count + 1}/2)")
            
            # Retry the task with exponential backoff
            raise self.retry(countdown=delay, exc=Exception(error_msg))
        
        # Update the guest translation record with actual character count
        try:
            latest_translation = GuestTranslation.query.filter_by(ip_address=client_ip).order_by(GuestTranslation.created_at.desc()).first()
            if latest_translation:
                latest_translation.character_count = character_count
                db.session.commit()
                print(f"Celery guest task {self.request.id}: Translation recorded for IP {client_ip}")
        except Exception as e:
            db.session.rollback()
            print(f"Celery guest task {self.request.id}: Error recording translation for IP {client_ip}: {e}")

        print(f"Celery guest task {self.request.id}: Translation successful ({src_lang} → {dest_lang}). File at: {translated_file_path}, Chars: {character_count}")
        
        # Try to upload to S3 with OSS fallback
        upload_result = {'success': False, 'service': None, 'key': None, 'download_url': None}
        
        if s3_service.is_available():
            try:
                upload_result = s3_service.upload_file_with_fallback(
                    translated_file_path, 
                    self.request.id, 
                    f"translated_{original_filename}"
                )
                
                if upload_result['success']:
                    service = upload_result['service']
                    key = upload_result['key']
                    download_url = upload_result['download_url']
                    
                    print(f"Celery guest task {self.request.id}: File uploaded to {service.upper()}: {key}")
                    
                    # Clean up local file after successful upload
                    try:
                        os.unlink(translated_file_path)
                        print(f"Celery guest task {self.request.id}: Local file cleaned up")
                    except Exception as e:
                        print(f"Celery guest task {self.request.id}: Warning - could not clean up local file: {e}")
                else:
                    print(f"Celery guest task {self.request.id}: Upload failed: {upload_result['error']}")
                    
            except Exception as e:
                print(f"Celery guest task {self.request.id}: Error during upload with fallback: {e}")
        else:
            print(f"Celery guest task {self.request.id}: S3 not available, using local storage")
        
        # Return information that the client can use
        result = {
            'status': 'SUCCESS',
            'message': 'File translated successfully.',
            'original_filename': original_filename,
            'character_count': character_count,
            'translation_rate': translation_rate,
            'texts_translated': int(translation_rate * len(original_texts)),
            'total_texts': len(original_texts)
        }
        
        # Include appropriate file reference based on storage method
        if upload_result['success']:
            result['storage_key'] = upload_result['key']
            result['download_url'] = upload_result['download_url']
            result['storage_type'] = upload_result['service']
            # Keep backward compatibility
            if upload_result['service'] == 's3':
                result['s3_key'] = upload_result['key']
        else:
            result['translated_file_path'] = translated_file_path  # Path on the WORKER dyno
            result['storage_type'] = 'local'
        
        return result
    except Retry:
        # Re-raise Retry exceptions to let Celery handle them
        raise
    except Exception as e:
        print(f"Celery guest task {self.request.id}: Error during translation: {e}")
        import traceback
        traceback.print_exc()
        # Re-raising will mark task as FAILED 
        raise 