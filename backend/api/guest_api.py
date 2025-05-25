"""
API endpoints for handling guest operations.
"""

from flask import Blueprint, jsonify, request, send_file, make_response, redirect
from services.user_service import get_guest_status, check_guest_permission
from services.translate_service import translate_pptx
from services.tasks import process_guest_translation_task
from pptx import Presentation
from db.models import GuestTranslation, db
from services.s3_service import s3_service
import os

guest_bp = Blueprint('guest', __name__)

@guest_bp.route('/api/guest/status', methods=['GET'])
def guest_status():
    """Get the guest translation status for the current IP."""
    client_ip = request.remote_addr
    
    # Get guest status
    status = get_guest_status(client_ip)
    return jsonify(status), 200

@guest_bp.route('/guest-translate', methods=['POST'])
def guest_translate_pptx_endpoint():
    """
    Endpoint for guest users to translate a PowerPoint file without authentication.
    Limited to GUEST_TRANSLATION_LIMIT per IP address.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
        
    file = request.files['file']
    
    # Check if file has a name
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
        
    # Check file extension
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({
            'error': 'Invalid file format. Only .pptx files are supported. Please save your PowerPoint file in the .pptx format and try again.'
        }), 400
        
    src_lang = request.form.get('src_lang', 'zh')  # Default source language: Chinese
    dest_lang = request.form.get('dest_lang', 'en')  # Default target language: English
    
    print(f"Received guest file: {file.filename}")
    print(f"Source language: {src_lang}, Target language: {dest_lang}")
    
    # Get client IP
    client_ip = request.remote_addr
    
    try:
        # Estimate the character count before translating
        estimated_character_count = 0
        try:
            # Make a copy of the file stream
            file_copy = file.stream.read()
            file.stream.seek(0)  # Reset the file pointer for later use
            
            # Load the presentation to count characters
            prs = Presentation(file_copy)
            
            # Count characters in text shapes and tables
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Handle regular text shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        estimated_character_count += len(shape.text)
                    
                    # Handle tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    estimated_character_count += len(cell.text)
                                    
            print(f"Estimated character count: {estimated_character_count}")
            file.stream.seek(0)  # Reset file pointer again after counting
        except Exception as e:
            print(f"Error estimating character count: {e}")
            # If estimation fails, continue with character count of 0

        # Check guest permission before translating
        result = check_guest_permission(client_ip, file.filename, src_lang, dest_lang, estimated_character_count)
        # Handle the case where three values are returned (allowed, response_obj, status_code)
        if isinstance(result, tuple):
            if len(result) == 3:
                allowed, response_obj, status_code = result
                if not allowed:
                    return response_obj, status_code
            elif len(result) == 2:
                allowed, response_obj = result
                if not allowed:
                    return response_obj
        
        # Translate the PPTX file and get the output path and character count
        output_path, character_count = translate_pptx(file.stream, src_lang, dest_lang)
        
        # Update the guest translation record with actual character count
        latest_translation = GuestTranslation.query.filter_by(ip_address=client_ip).order_by(GuestTranslation.created_at.desc()).first()
        if latest_translation:
            latest_translation.character_count = character_count
            db.session.commit()
        
        # Return the translated file
        response = make_response(send_file(output_path, as_attachment=True, 
                             download_name=f"translated_{file.filename}"))
        
        # Set content type explicitly (helps prevent MIME type issues)
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        print(f"Guest translation completed successfully. Used {character_count} characters.")
        return response
    except Exception as e:
        print(f"Error during guest translation: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@guest_bp.route('/guest-translate-async-start', methods=['POST'])
def guest_translate_async_start_endpoint():
    """Endpoint to start asynchronous translation for guest users."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
        
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
        
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({
            'error': 'Invalid file format. Only .pptx files are supported.'
        }), 400
        
    src_lang = request.form.get('src_lang', 'zh')
    dest_lang = request.form.get('dest_lang', 'en')
    
    print(f"API: Received guest file for async translation: {file.filename}")
    print(f"API: Source lang: {src_lang}, Target lang: {dest_lang}")

    # Get client IP
    client_ip = request.remote_addr

    try:
        # Estimate the character count before translating
        estimated_character_count = 0
        try:
            # Make a copy of the file stream
            file_copy = file.stream.read()
            file.stream.seek(0)  # Reset the file pointer for later use
            
            # Load the presentation to count characters
            prs = Presentation(file_copy)
            
            # Count characters in text shapes and tables
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Handle regular text shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        estimated_character_count += len(shape.text)
                    
                    # Handle tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    estimated_character_count += len(cell.text)
                                    
            print(f"Estimated character count: {estimated_character_count}")
            file.stream.seek(0)  # Reset file pointer again after counting
        except Exception as e:
            print(f"Error estimating character count: {e}")
            # If estimation fails, continue with character count of 0

        # Check guest permission before translating
        result = check_guest_permission(client_ip, file.filename, src_lang, dest_lang, estimated_character_count)
        # Handle the case where three values are returned (allowed, response_obj, status_code)
        if isinstance(result, tuple):
            if len(result) == 3:
                allowed, response_obj, status_code = result
                if not allowed:
                    return response_obj, status_code
            elif len(result) == 2:
                allowed, response_obj = result
                if not allowed:
                    return response_obj

        # Read file content into bytes to pass to Celery task
        file_bytes = file.read()
        file.stream.seek(0) # Reset stream position

        # Dispatch the Celery task for guest translation
        task = process_guest_translation_task.delay(
            client_ip=client_ip,
            original_file_content_bytes=file_bytes, 
            original_filename=file.filename, 
            src_lang=src_lang, 
            dest_lang=dest_lang,
            estimated_character_count=estimated_character_count
        )
        
        print(f"API: Dispatched guest Celery task ID: {task.id}")
        return jsonify({'message': 'Translation task started', 'task_id': task.id}), 202
    except Exception as e:
        print(f"API: Error dispatching guest Celery task: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': 'Could not start translation task', 'details': str(e)}), 500

@guest_bp.route('/guest-translate-status/<task_id>', methods=['GET'])
def get_guest_translation_status_endpoint(task_id):
    """Endpoint to check the status of a guest translation task."""
    task = process_guest_translation_task.AsyncResult(task_id)
    
    response_data = {
        'task_id': task_id,
        'status': task.state
    }
    
    if task.state == 'PENDING':
        response_data['message'] = 'Task is pending or unknown.'
    elif task.state == 'PROGRESS':
        response_data['message'] = 'Task is in progress.'
        response_data['progress'] = task.info.get('progress', 0) 
    elif task.state == 'SUCCESS':
        response_data['message'] = 'Task completed successfully.'
        response_data['result'] = task.result 
        # Add download URL for the translated file
        if task.result:
            if task.result.get('storage_type') == 's3' and task.result.get('download_url'):
                # For S3 storage, use the presigned URL directly
                response_data['result']['download_url'] = task.result['download_url']
            elif 'translated_file_path' in task.result:
                # For local storage, use the download endpoint
                response_data['result']['download_url'] = f'/guest-download/{task_id}'
    elif task.state == 'FAILURE':
        response_data['message'] = 'Task failed.'
        response_data['error'] = 'An error occurred during translation processing.' 
    
    return jsonify(response_data)

@guest_bp.route('/guest-download/<task_id>', methods=['GET'])
def download_guest_translated_file(task_id):
    """Endpoint to download the translated file for guest users."""
    task = process_guest_translation_task.AsyncResult(task_id)
    
    if task.state != 'SUCCESS':
        return jsonify({'error': 'Translation not completed or failed'}), 400
        
    if not task.result:
        return jsonify({'error': 'Translated file not found'}), 404
    
    # Handle S3 storage
    if task.result.get('storage_type') == 's3':
        if task.result.get('download_url'):
            # Redirect to the S3 presigned URL
            return redirect(task.result['download_url'])
        elif task.result.get('s3_key'):
            # Generate a new presigned URL if the old one expired
            download_url = s3_service.generate_presigned_url(task.result['s3_key'], expiration=3600)
            if download_url:
                return redirect(download_url)
            else:
                return jsonify({'error': 'Could not generate download URL'}), 500
        else:
            return jsonify({'error': 'S3 file reference not found'}), 404
    
    # Handle local storage (fallback)
    elif 'translated_file_path' in task.result:
        file_path = task.result['translated_file_path']
        original_filename = task.result.get('original_filename', 'translated_file.pptx')
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File no longer exists'}), 404
            
        try:
            response = make_response(send_file(file_path, as_attachment=True, 
                                             download_name=f"translated_{original_filename}"))
            response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            return response
        except Exception as e:
            print(f"Error downloading guest file: {e}")
            return jsonify({'error': 'Error downloading file'}), 500
    else:
        return jsonify({'error': 'No file reference found'}), 404 