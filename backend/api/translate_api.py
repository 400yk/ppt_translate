"""
API endpoints for handling PowerPoint translation requests.
"""
from flask import request, send_file, jsonify, make_response, Blueprint
from flask_jwt_extended import jwt_required, get_jwt_identity
from db.models import User, db
from services.user_service import check_user_permission, check_guest_permission
from services.translate_service import translate_pptx
from pptx import Presentation

translate_bp = Blueprint('translate', __name__)

@translate_bp.route('/translate', methods=['POST'])
@jwt_required()
def translate_pptx_endpoint():
    """Endpoint to translate PowerPoint files."""
    username = get_jwt_identity()
    user = User.query.filter_by(username=username).first()
    
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
        
    file = request.files['file']
    
    # Check if file has a name
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
        
    # Check file extension
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({
            'error': 'Invalid file format. Only .pptx files are supported. Please save your PowerPoint file in the .pptx format and try again.'
        }), 400
        
    src_lang = request.form.get('src_lang', 'zh')  # Default source language: Chinese
    dest_lang = request.form.get('dest_lang', 'en')  # Default target language: English
    
    print(f"Received file: {file.filename}")
    print(f"Source language: {src_lang}, Target language: {dest_lang}")
    
    try:
        # Check user's permission to translate - MOVED BEFORE translation
        result = check_user_permission(user)
        # Handle the case where three values are returned (allowed, response_obj, status_code)
        if isinstance(result, tuple):
            if len(result) == 3:
                allowed, response_obj, status_code = result
                if not allowed:
                    return response_obj, status_code
            elif len(result) == 2:
                allowed, response_obj = result
                if not allowed:
                    return response_obj
        
        # Translate the PPTX file and get the output path and character count
        output_path, character_count = translate_pptx(file.stream, src_lang, dest_lang)
        
        # Record the translation in the database including character count
        user.record_translation(file.filename, src_lang, dest_lang, character_count)
        
        # Return the translated file
        response = make_response(send_file(output_path, as_attachment=True, 
                             download_name=f"translated_{file.filename}"))
        
        # Set content type explicitly (helps prevent MIME type issues)
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        print(f"Translation completed successfully. Used {character_count} characters.")
        return response
    except Exception as e:
        print(f"Error during translation: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@translate_bp.route('/guest-translate', methods=['POST'])
def guest_translate_pptx_endpoint():
    """
    Endpoint for guest users to translate a PowerPoint file without authentication.
    Limited to GUEST_TRANSLATION_LIMIT per IP address.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
        
    file = request.files['file']
    
    # Check if file has a name
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
        
    # Check file extension
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({
            'error': 'Invalid file format. Only .pptx files are supported. Please save your PowerPoint file in the .pptx format and try again.'
        }), 400
        
    src_lang = request.form.get('src_lang', 'zh')  # Default source language: Chinese
    dest_lang = request.form.get('dest_lang', 'en')  # Default target language: English
    
    print(f"Received guest file: {file.filename}")
    print(f"Source language: {src_lang}, Target language: {dest_lang}")
    
    # Get client IP
    client_ip = request.remote_addr
    
    try:
        # Estimate the character count before translating
        estimated_character_count = 0
        try:
            # Make a copy of the file stream
            file_copy = file.stream.read()
            file.stream.seek(0)  # Reset the file pointer for later use
            
            # Load the presentation to count characters
            prs = Presentation(file_copy)
            
            # Count characters in text shapes and tables
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Handle regular text shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        estimated_character_count += len(shape.text)
                    
                    # Handle tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    estimated_character_count += len(cell.text)
                                    
            print(f"Estimated character count: {estimated_character_count}")
            file.stream.seek(0)  # Reset file pointer again after counting
        except Exception as e:
            print(f"Error estimating character count: {e}")
            # If estimation fails, continue with character count of 0

        # Check guest permission before translating
        result = check_guest_permission(client_ip, file.filename, src_lang, dest_lang, estimated_character_count)
        # Handle the case where three values are returned (allowed, response_obj, status_code)
        if isinstance(result, tuple):
            if len(result) == 3:
                allowed, response_obj, status_code = result
                if not allowed:
                    return response_obj, status_code
            elif len(result) == 2:
                allowed, response_obj = result
                if not allowed:
                    return response_obj
        
        # Translate the PPTX file and get the output path and character count
        output_path, character_count = translate_pptx(file.stream, src_lang, dest_lang)
        
        # Return the translated file
        response = make_response(send_file(output_path, as_attachment=True, 
                             download_name=f"translated_{file.filename}"))
        
        # Set content type explicitly (helps prevent MIME type issues)
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        print(f"Guest translation completed successfully. Used {character_count} characters.")
        return response
    except Exception as e:
        print(f"Error during guest translation: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500
        
@translate_bp.route('/api/translations/history', methods=['GET'])
@jwt_required()
def get_translation_history():
    username = get_jwt_identity()
    user = User.query.filter_by(username=username).first()
    
    if not user:
        return jsonify({'error': 'User not found'}), 404
        
    # Get the user's translation history
    translations = user.translations.order_by(db.desc('created_at')).all()
    
    history = []
    for translation in translations:
        history.append({
            'id': translation.id,
            'filename': translation.filename,
            'date': translation.created_at.isoformat(),
            'source_language': translation.source_language,
            'target_language': translation.target_language
        })
        
    return jsonify({
        'count': len(history),
        'translations': history
    }), 200 