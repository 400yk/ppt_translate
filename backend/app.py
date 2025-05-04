import os
import tempfile
from flask import Flask, request, send_file, jsonify
from pptx import Presentation
import requests
import ast  # Added for safe eval fallback
from dotenv import load_dotenv
from pptx.util import Pt
from pptx_utils import measure_text_bbox, fit_font_size_to_bbox, fit_font_size_for_title
from config import DEFAULT_FONT_NAME, DEFAULT_FONT_SIZE, GEMINI_API_URL, DEFAULT_TITLE_FONT_SIZE
from pptx.enum.shapes import PP_PLACEHOLDER

# Load Gemini API key from .env - try multiple locations
current_dir = os.path.dirname(os.path.abspath(__file__))
root_dir = os.path.dirname(current_dir)

# Try loading .env from different possible locations
env_paths = [
    os.path.join(root_dir, '.env'),  # root project directory
    os.path.join(current_dir, '.env'),  # backend directory
    '.env'  # current working directory
]

for env_path in env_paths:
    if os.path.exists(env_path):
        print(f"Loading environment from: {env_path}")
        load_dotenv(env_path)
        break
else:
    print("Warning: Could not find .env file in any of the expected locations")

GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
if not GEMINI_API_KEY:
    print("ERROR: GEMINI_API_KEY not found in environment variables")

from flask_cors import CORS
app = Flask(__name__)
CORS(app)

def gemini_batch_translate(texts, src_lang, dest_lang):
    """Batch translate a list of texts using Gemini API."""
    # Check if API key is available
    if not GEMINI_API_KEY:
        print("ERROR: Missing Gemini API key. Translation will return original text.")
        return texts
        
    # Join texts as a JSON list to preserve order and mapping
    import json
    joined = json.dumps(texts, ensure_ascii=False)
    prompt = f"Translate the following JSON list from {src_lang} to {dest_lang}. Only return the translated JSON list, no explanations.\n{joined}"
    headers = {
        'Content-Type': 'application/json',
    }
    params = {
        'key': GEMINI_API_KEY
    }
    data = {
        'contents': [{
            'parts': [{
                # Instruct Gemini to return a strict JSON array
                'text': prompt + "\nReturn the result strictly as a JSON array (not a Python list), no extra text."
            }]
        }]
    }
    try:
        resp = requests.post(GEMINI_API_URL, headers=headers, params=params, json=data, timeout=60)
        resp.raise_for_status()
        result = resp.json()
        if 'candidates' in result and result['candidates']:
            translated_json = result['candidates'][0]['content']['parts'][0]['text']
            # Strip code block markers from Gemini response
            if translated_json.strip().startswith('```'):
                translated_json = translated_json.strip().lstrip('`').lstrip('json').strip()
                if translated_json.endswith('```'):
                    translated_json = translated_json[:translated_json.rfind('```')].strip()
            # Try to parse the JSON list from Gemini's response
            try:
                return json.loads(translated_json)
            except Exception as e:
                print(f"Error parsing Gemini JSON: {e}\nRaw: {translated_json}")
                # Try ast.literal_eval as a fallback for Python-style lists
                try:
                    return ast.literal_eval(translated_json)
                except Exception as e2:
                    print(f"Fallback ast.literal_eval failed: {e2}\nRaw: {translated_json}")
                    return texts  # fallback
        elif 'error' in result:
            error_msg = result.get('error', {}).get('message', 'Unknown error')
            print(f"Gemini API error: {error_msg}")
            
            # Special handling for auth errors
            if 'API key not valid' in error_msg or '403' in error_msg:
                print("Authentication error: Please check your Gemini API key.")
            
            raise Exception(f"Gemini API error: {error_msg}")
        else:
            raise Exception("Unexpected Gemini API response")
    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 403:
            print(f"Gemini API authentication error (403 Forbidden): Please check your API key.")
            print(f"API Response: {e.response.text}")
        else:
            print(f"Gemini HTTP error: {e}")
        return texts  # fallback to original
    except Exception as e:
        print(f"Gemini translation error: {e}")
        return texts  # fallback to original

def translate_pptx(input_stream, src_lang, dest_lang):
    prs = Presentation(input_stream)
    text_shapes = []
    texts = []
    table_cells = []
    table_texts = []
    
    # Collect all text shapes and their texts
    for slide in prs.slides:
        for shape in slide.shapes:
            # Handle regular text shapes
            if hasattr(shape, "text") and shape.text.strip():
                text_shapes.append(shape)
                texts.append(shape.text)
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            table_cells.append(cell)
                            table_texts.append(cell.text)

    # Batch translate all text content together
    all_texts = texts + table_texts
    all_translated_texts = gemini_batch_translate(all_texts, src_lang, dest_lang)
    
    # Split the translated texts back into shape texts and table texts
    translated_texts = all_translated_texts[:len(texts)]
    translated_table_texts = all_translated_texts[len(texts):]
    
    print(f"Collected {len(texts)} texts from regular shapes.")
    print(f"Collected {len(table_texts)} texts from table cells.")
    print(f"Received {len(all_translated_texts)} translated texts total.")
    print(f"Sample original: {all_texts[:3]}")
    print(f"Sample translated: {all_translated_texts[:3]}")
    
    # Update regular text shapes
    for shape, translated in zip(text_shapes, translated_texts):
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            text_frame = shape.text_frame
            # Get original font properties from the first run (if available)
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                original_run = text_frame.paragraphs[0].runs[0]
                font_name = original_run.font.name or DEFAULT_FONT_NAME
                
                # Style scale factor for text formatting (bold/italic/underline)
                style_scale_factor = 1.0
                
                # Check if text has formatting that might need more space
                if original_run.font.bold:
                    style_scale_factor *= 0.9  # Bold text needs ~10% more space
                if original_run.font.italic:
                    style_scale_factor *= 0.95  # Italic text needs ~5% more space
                if original_run.font.underline:
                    style_scale_factor *= 0.95  # Underlined text needs ~5% more space
                
                # Check if this is a title placeholder
                is_title = False
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Title placeholders have types: TITLE (1) or CENTER_TITLE (3)
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        is_title = True
                
                # Use appropriate default font size based on whether it's a title
                if original_run.font.size:
                    original_font_size = int(original_run.font.size.pt)
                else:
                    original_font_size = DEFAULT_TITLE_FONT_SIZE if is_title else DEFAULT_FONT_SIZE
                
                # Enhanced font color extraction: run -> paragraph -> text_frame
                font_color = None
                # 1. Try run-level color
                if original_run.font.color is not None:
                    try:
                        font_color = original_run.font.color.rgb
                    except AttributeError:
                        font_color = None
                # 2. Try paragraph-level color
                if not font_color:
                    para = text_frame.paragraphs[0]
                    if para.font and para.font.color is not None:
                        try:
                            font_color = para.font.color.rgb
                        except AttributeError:
                            font_color = None
                # 3. Try text_frame-level color
                if not font_color and hasattr(text_frame, 'font') and text_frame.font is not None:
                    try:
                        font_color = text_frame.font.color.rgb
                    except AttributeError:
                        font_color = None
            else:
                font_name = DEFAULT_FONT_NAME
                style_scale_factor = 1.0  # Default for no formatting
                
                # Check if this is a title placeholder
                is_title = False
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        is_title = True
                
                original_font_size = DEFAULT_TITLE_FONT_SIZE if is_title else DEFAULT_FONT_SIZE
                font_color = None
                
            # Measure the bounding box of the original text
            original_text = shape.text
            orig_w, orig_h = measure_text_bbox(original_text, font_name, original_font_size)
            
            # For titles, only constrain height, not width
            if is_title:
                best_font_size = fit_font_size_for_title(orig_h, translated, font_name, original_font_size)
            else:
                # For regular content, constrain both dimensions
                best_font_size = fit_font_size_to_bbox(orig_w, orig_h, translated, font_name, original_font_size)
            
            # Apply the style scaling factor to the font size
            best_font_size = int(best_font_size * style_scale_factor)
                
            # Store original formatting information before clearing
            original_paragraphs = []
            for para in text_frame.paragraphs:
                para_info = {
                    'level': para.level,
                    'alignment': para.alignment,
                    'runs': []
                }
                
                # Store runs and their formatting
                for run in para.runs:
                    run_info = {
                        'text_length': len(run.text),
                        'bold': run.font.bold,
                        'italic': run.font.italic, 
                        'underline': run.font.underline,
                        'font_name': run.font.name or font_name,
                        'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') else None
                    }
                    para_info['runs'].append(run_info)
                
                original_paragraphs.append(para_info)
                
            text_frame.clear()
            
            # Simple approach for short texts - use a single formatted run
            if len(original_paragraphs) == 1 and len(original_paragraphs[0]['runs']) == 1:
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = translated
                run.font.name = font_name
                run.font.size = Pt(best_font_size)
                
                # Apply original formatting
                original_run_info = original_paragraphs[0]['runs'][0]
                if original_run_info['bold'] is not None:
                    run.font.bold = original_run_info['bold']
                if original_run_info['italic'] is not None:
                    run.font.italic = original_run_info['italic']
                if original_run_info['underline'] is not None:
                    run.font.underline = original_run_info['underline']
                if font_color:
                    run.font.color.rgb = font_color
            else:
                # For more complex text, preserve the structure using paragraph count
                # This is a simplification as we can't perfectly map formatting from source to translated text
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = translated
                run.font.name = font_name
                run.font.size = Pt(best_font_size)
                
                # Apply formatting from dominant formatting in original text
                has_bold = any(run_info['bold'] for para in original_paragraphs for run_info in para['runs'] if run_info['bold'] is not None)
                has_italic = any(run_info['italic'] for para in original_paragraphs for run_info in para['runs'] if run_info['italic'] is not None)
                has_underline = any(run_info['underline'] for para in original_paragraphs for run_info in para['runs'] if run_info['underline'] is not None)
                
                run.font.bold = has_bold
                run.font.italic = has_italic
                run.font.underline = has_underline
                
                if font_color:
                    run.font.color.rgb = font_color
        else:
            shape.text = translated
    
    # Update table cells with translated text
    for cell, translated in zip(table_cells, translated_table_texts):
        text_frame = cell.text_frame
        
        # Store original formatting if available
        original_paragraphs = []
        font_name = DEFAULT_FONT_NAME
        font_color = None
        original_font_size = DEFAULT_FONT_SIZE
        style_scale_factor = 1.0
        
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            original_run = text_frame.paragraphs[0].runs[0]
            font_name = original_run.font.name or DEFAULT_FONT_NAME
            
            # Style scale factor for text formatting (bold/italic/underline)
            style_scale_factor = 1.0
            
            # Check if text has formatting that might need more space
            if original_run.font.bold:
                style_scale_factor *= 0.9  # Bold text needs ~10% more space
            if original_run.font.italic:
                style_scale_factor *= 0.95  # Italic text needs ~5% more space
            if original_run.font.underline:
                style_scale_factor *= 0.95  # Underlined text needs ~5% more space
            
            # Get original font size
            if original_run.font.size:
                original_font_size = int(original_run.font.size.pt)
                
            # Try to get font color
            if original_run.font.color is not None:
                try:
                    font_color = original_run.font.color.rgb
                except AttributeError:
                    font_color = None
                    
            # Store original formatting information
            for para in text_frame.paragraphs:
                para_info = {
                    'level': para.level,
                    'alignment': para.alignment,
                    'runs': []
                }
                
                # Store runs and their formatting
                for run in para.runs:
                    run_info = {
                        'text_length': len(run.text),
                        'bold': run.font.bold,
                        'italic': run.font.italic, 
                        'underline': run.font.underline,
                        'font_name': run.font.name or font_name,
                        'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') else None
                    }
                    para_info['runs'].append(run_info)
                
                original_paragraphs.append(para_info)
        
        # Measure the bounding box of the original text - for table cells
        # Table cells have fixed width, so we need to fit text within that constraint
        original_text = cell.text
        
        # Get the approximate cell dimensions
        # For table cells, we calculate width and height based on the cell's margin and text area
        # We use the text_frame margins to estimate the available space
        margin_left = text_frame.margin_left.inches if hasattr(text_frame, 'margin_left') else 0.1
        margin_right = text_frame.margin_right.inches if hasattr(text_frame, 'margin_right') else 0.1
        margin_top = text_frame.margin_top.inches if hasattr(text_frame, 'margin_top') else 0.05
        margin_bottom = text_frame.margin_bottom.inches if hasattr(text_frame, 'margin_bottom') else 0.05
        
        # Estimate the available width and height in the cell
        cell_width = (cell.width.inches if hasattr(cell, 'width') else 2.0) - margin_left - margin_right
        cell_height = (cell.height.inches if hasattr(cell, 'height') else 0.5) - margin_top - margin_bottom
        
        # Convert to points for measurement
        cell_width_pt = cell_width * 72  # 72 points per inch
        cell_height_pt = cell_height * 72
        
        # Measure the original text bounding box
        orig_w, orig_h = measure_text_bbox(original_text, font_name, original_font_size)
        
        # Calculate the best font size to fit the translated text in the cell
        best_font_size = fit_font_size_to_bbox(cell_width_pt, cell_height_pt, translated, font_name, original_font_size)
        
        # Apply the style scaling factor to the font size
        best_font_size = int(best_font_size * style_scale_factor)
        
        # Clear the text frame and add the translated text
        text_frame.clear()
        
        # Simple approach for table cells - use a single run to maintain consistent formatting
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = translated
        run.font.name = font_name
        run.font.size = Pt(best_font_size)
        
        # Apply original formatting
        if original_paragraphs:
            # Determine dominant formatting across all runs
            has_bold = any(run_info['bold'] for para in original_paragraphs 
                          for run_info in para['runs'] if run_info['bold'] is not None)
            has_italic = any(run_info['italic'] for para in original_paragraphs 
                            for run_info in para['runs'] if run_info['italic'] is not None)
            has_underline = any(run_info['underline'] for para in original_paragraphs 
                               for run_info in para['runs'] if run_info['underline'] is not None)
            
            run.font.bold = has_bold
            run.font.italic = has_italic
            run.font.underline = has_underline
        
        # Apply font color if available
        if font_color:
            run.font.color.rgb = font_color
            
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    temp_file.close()
    return temp_file.name

if __name__ == '__main__':
    app.run(debug=True)
